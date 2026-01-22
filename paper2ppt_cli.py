#!/usr/bin/env python3
"""
Paper2PPT CLI - Core Logic

This module handles the extraction of content from PDF papers and the structural 
generation of PowerPoint slides. It includes logic for:
- PDF Text extraction and cleaning
- Semantic section identification (Method, Results, etc.)
- Bullet point generation and refinement
- Image extraction and placement
"""
import warnings
warnings.filterwarnings("ignore")

import argparse
import os
import re
from pathlib import Path
from typing import List
from pptx import Presentation # type: ignore

from paper2ppt_core.io import load_input_paper
from paper2ppt_core.sections import split_into_sections
from paper2ppt_core.pptx_builder import build_presentation, MAX_FIGURES_PER_SLIDE
from ppt_narration_project.narration_generator import generate_narration


# pyright: reportMissingImports=false
from pptx import Presentation


# ==============================
# Optional LLM (SAFE OFF by default)
# ==============================
try:
    from models.qwen_llm import qwen_generate  # type: ignore
    HAS_LLM = True
except Exception:
    HAS_LLM = False
    def qwen_generate(prompt: str, max_tokens: int = 64, temperature: float = 0.1) -> str:
        return ""

# ==============================
# CONSTANTS (CONSERVATIVE)
# ==============================
MAX_MODEL_CHARS = 12000
MAX_BULLET_WORDS = 60
MIN_BULLET_WORDS = 6

MAX_BULLETS_PER_SLIDE = 6
MIN_BULLETS_PER_SLIDE = 3



SKIP_SECTIONS = {
    "references", "acknowledgements",
    "appendix", "supplementary"
}

SECTION_MAP = {
    "abstract": "Overview",
    "introduction": "Introduction",
    "background": "Background",
    "related work": "Related Work",
    "method": "Method",
    "model": "Method",
    "architecture": "Method",
    "approach": "Method",
    "experiment": "Experiments",
    "evaluation": "Results",
    "result": "Results",
    "discussion": "Discussion",
    "conclusion": "Conclusion",
}

SECTION_TARGET_BULLETS = {
    "Overview": 4,
    "Introduction": 6,
    "Background": 4,
    "Related Work": 4,
    "Method": 8,
    "Results": 6,
    "Conclusion": 4,
}

# ==============================
# TEXT NORMALIZATION
# ==============================
def normalize_pdf_text(text: str) -> str:
    text = re.sub(r'\n?\d+(\.\d+)*\s+[A-Z][A-Za-z\s\-]{3,}', ' ', text)
    text = text.replace("\n", " ")
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

def normalize_section(title: str) -> str:
    t = title.lower()
    for k, v in SECTION_MAP.items():
        if k in t:
            return v
    return title.title()

# ==============================
# SENTENCE EXTRACTION (SAFE)
# ==============================
def extract_sentences(text: str) -> List[str]:
    """
    Extract full, grammatical sentences only.
    """
    sentences = re.split(r'(?<=[.!?])\s+', text)
    out = []
    for s in sentences:
        wc = len(s.split())
        if 8 <= wc <= 80:
            out.append(s.strip())
    return out

# ==============================
# BULLET GENERATION (NON-DESTRUCTIVE)
# ==============================
def rewrite_bullet(sentence: str) -> str:

     # --- FIX PDF HYPHENATION (GLOBAL) ---
    sentence = re.sub(r'(\w)-\s+(\w)', r'\1\2', sentence)

    s = re.sub(r"\s+", " ", sentence.strip())
    
    # --- CLEANUP PREFIXES (New) ---
    # Remove "However,", "Thus,", "Therefore,", "In this paper," etc.
    s = re.sub(r'^(however|therefore|thus|moreover|furthermore|consequently|hence|accordingly|specifically|notably|importantly|interestingly|finally|additionally)[, ]+', '', s, flags=re.IGNORECASE)
    s = re.sub(r'^(in this paper|in this work|we show that|we demonstrate that|we find that|it is observed that)[, ]+', '', s, flags=re.IGNORECASE)
    
    # Capitalize after strip
    if s:
        s = s[0].upper() + s[1:]

    low = s.lower()

    # --- GLOBAL STRUCTURAL REJECTIONS ---

    # 1. Reject fragments starting with conjunctions
    if re.match(r'^(and|but|or)\b', low):
        return ""

    # 2. Reject citation residue / orphaned parentheses
    if re.match(r'^\)?\d{4}\)', s) or s.startswith(")"):
        return ""

    # 3. Reject ellipsis fragments (not full sentences)
    if "…" in s or "..." in s:
        return ""

    # Reject numeric-only figure references (e.g., "5 illustrates ...")
    if re.match(r'^\d+\s+(illustrates|shows|presents)', low):
        return ""


    # 2. Reject broken comparative or numeric claims
    if re.search(r"\b(by over|achieves|improves|outperforms)\b\s*(,|\.|$)", low):
        return ""

    # 3. Reject dangling prepositions at sentence end
    if re.search(r"\b(by|over|of|with|to|for|in|on)\s*\.$", low):
        return ""

    # 4. Reject purely navigational table / figure / section references
    # (Revised to ALLOW "Table 1 shows..." but REJECT "See Table 1")
    if re.search(r"^\s*(see|refer to|shown in|details in|as seen in)\s+(table|figure|fig\.|section)\s+\d+", low):
        return ""
    # Reject short captions disguised as sentences
    if re.match(r"^(table|figure|fig\.)\s+\d+\.?\s*$", low):
        return ""

    # 5. Reject obvious author / repo / citation noise
    if any(x in low for x in [
        "et al.", "arxiv", "github", "codebase",
        "implemented by", "designed by"
    ]):
        return ""

    if not s:
        return ""

    # Normalize whitespace
    s = re.sub(r'\s+', ' ', s.strip())

    low = s.lower()

    # Drop obvious noise / metadata
    if any(x in low for x in [
        "et al.", "arxiv", "acknowledgement", "references",
        "conference", "proceedings",
        "nips", "neurips"
        # Removed "table", "figure" from here to allow legitimate discussion
    ]):
        return ""

    # Remove citation brackets safely (does NOT cut sentence)
    s = re.sub(r'\([^)]*\)', '', s)
    s = re.sub(r'\[[^\]]*\]', '', s)
    s = re.sub(r'\s+', ' ', s).strip()

    words = s.split()

    # Length guard (conservative)
    # Using MAX_BULLET_WORDS constant (60)
    if len(words) < MIN_BULLET_WORDS or len(words) > MAX_BULLET_WORDS:
        return ""

    # Capitalize safely
    if s:
        s = s[0].upper() + s[1:]

    # Ensure proper sentence ending
    if not s.endswith("."):
        s += "."

    return s

# ==============================

def is_complete_sentence(sentence: str) -> bool:
    """
    Rejects bullets that are syntactically finished
    but semantically incomplete.
    """
    s = sentence.lower().strip()

    # Must end with a period
    if not s.endswith("."):
        return False

    # Must contain a verb
    if not re.search(
        r"\b(is|are|was|were|has|have|achieves|uses|shows|demonstrates|improves|reduces|introduces|presents|validates|evaluates|employs|contains|includes|consists|resulted|outperformed)\b",
        s
    ):
        return False

    # Reject dangling endings
    BAD_ENDINGS = (
        "and.", "or.", "with.", "by.", "to.", "that.", "which.",
        "including.", "based on.", "consisting of."
    )
    for end in BAD_ENDINGS:
        if s.endswith(end):
            return False

    # Reject obvious truncation tokens
    if any(x in s for x in ["…", "...", " d…", " de…"]):
        return False

    return True


def polish_bullets(bullets):
    bullets = [b for b in bullets if "…" not in b and "..." not in b]
    final = []
    for b in bullets:
        text = b[2:] if b.startswith("* ") else b
        text = text.rstrip(",;:")
        if not text.endswith("."):
            text += "."
        final.append(f"* {text}")
    return final

def finalize_bullet(bullet: str) -> str:
    """
    Final defensive fix for broken academic sentences.
    Does NOT change meaning.
    Does NOT remove bullets.
    """

    text = bullet.strip()
    if text.startswith("* "):
        text = text[2:].strip()

    text = text.rstrip(",;:")

    # Detect obvious truncation artifacts
    if re.search(r'\.\.\.$', text):
        text = text.replace("...", " in prior work.")

    if re.search(r'\b(convo|transduct|averagin|repr|competitiv|performan)\b$', text.lower()):
        text += " settings."

    if re.search(r'\b(achieves|improving|by over)\b\s*$', text.lower()):
        text += " existing baselines."

    if re.search(r'\b(of any of the)\b$', text.lower()):
        text += " compared approaches."

    # Final period
    if not text.endswith("."):
        text += "."

    # Capitalize safely
    text = text[0].upper() + text[1:]

    return f"* {text}"



# ==============================
# BULLET FILTER (VERY LIGHT)
# ==============================
def final_bullets(candidates: List[str]) -> List[str]:
    """
    Final refinement step for bullet points.
    - Accepts valid academic sentences
    - Prevents duplicates and fragments
    - Ensures punctuation
    """
    out = []
    seen = set()

    for b in candidates:
        if not b:
            continue

        text = b.strip()
        if text.startswith("* "):
            text = text[2:].strip()

        words = text.split()
        wc = len(words)

        # Allow slightly shorter bullets than before (content boost)
        if wc < MIN_BULLET_WORDS:
            continue
        if wc > MAX_BULLET_WORDS:
            continue

        low = text.lower()

        # Must have meaning signal
        has_signal = (
            re.search(r"\b(is|are|was|were|uses|achieves|shows|demonstrates|improves|reduces|introduces|presents|stacks|computes|trains|learns|optimizes|functions|generates|outputs|inputs|consists|comprises|employs|utilizes|applies|contains|includes)\b", low)
            or re.search(r"\b(model|approach|method|architecture|framework|system|mechanism|network|layer|attention|encoder|decoder|data|training|loss|transformer|embedding|projection|softmax|normalization)\b", low)
        )

        if not has_signal:
            continue

        # Deduplicate aggressively
        key = re.sub(r"\W+", "", low)
        if key in seen:
            continue
        seen.add(key)

        # Final safety ending
        text = text.rstrip(",;:")
        if not text.endswith("."):
            text += "."

        text = text[0].upper() + text[1:]
        out.append(f"* {text}")

    return out

# ==============================
# SECURE EXPANSION (SAFE)
# ==============================
def add_one_more_safe_bullet(
    existing_bullets: List[str],
    all_sentences: List[str],
    target_count: int
) -> List[str]:
    """
    Adds at most ONE extra bullet if section is thin.
    Uses only already-safe, complete sentences.
    """

    if len(existing_bullets) >= target_count:
        return existing_bullets

    used_text = {b.lower() for b in existing_bullets}

    for s in all_sentences:
        s = s.strip()

        # Must be full sentence
        if not s.endswith("."):
            continue

        wc = len(s.split())
        if wc < MIN_BULLET_WORDS or wc > 45:
            continue

        low = s.lower()

        # Reject dangling endings
        if re.search(r'\b(and|or|with|that|which|to|by)\.$', low):
            continue

        # Must contain a verb-like signal
        if not re.search(
            r"\b(is|are|was|were|uses|shows|demonstrates|achieves|improves|reduces|introduces|presents|stacks|computes|trains|learns|utilizes|employs|generates)\b",
            low,
        ):
            continue

        # Avoid duplicates
        if low in used_text:
            continue

        # SAFE to add
        rb = rewrite_bullet(s)
        if not rb:
            continue

        return existing_bullets + [f"* {rb}"]


    return existing_bullets



# ==============================
# STEP-3: LIGHT BULLET POLISH (SAFE)
# ==============================

_BAD_ENDINGS = (
    "and", "or", "with", "by", "to", "that", "which", "including", "based on"
)

def light_polish_bullet(bullet: str) -> str:
    """
    Very light, non-destructive polish.
    Fixes form only, not meaning.
    """
    b = bullet.strip()

    # Remove "* " temporarily
    if b.startswith("* "):
        b = b[2:].strip()

    low = b.lower()

    # Fix dangling endings
    for end in _BAD_ENDINGS:
        if low.endswith(" " + end) or low.endswith(" " + end + "."):
            b = b.rstrip(".") + " in practice."

    # Remove trailing commas / colons
    b = b.rstrip(",;:")

    # Capitalize first letter
    b = b[0].upper() + b[1:]

    # Ensure sentence ends properly
    if not b.endswith("."):
        b += "."

    return f"* {b}"


# ==============================
# FLOW + CHUNKING
# ==============================
def enhance_bullet_flow(bullets: List[str]) -> List[str]:
    """
    Simple ordering only. Never removes bullets.
    """
    raw = [b[2:] if b.startswith("* ") else b for b in bullets]

    priority = {
        "problem": 0,
        "challenge": 0,
        "approach": 1,
        "method": 1,
        "model": 1,
        "result": 2,
        "performance": 2,
        "conclusion": 3
    }

    def score(t: str) -> int:
        tl = t.lower()
        for k, v in priority.items():
            if k in tl:
                return v
        return 1

    raw.sort(key=score)
    return [f"* {b}" for b in raw]

def chunk_bullets(bullets: List[str]) -> List[List[str]]:
    chunks, i = [], 0
    while i < len(bullets):
        chunk = bullets[i:i + MAX_BULLETS_PER_SLIDE]
        if len(chunk) < MIN_BULLETS_PER_SLIDE and chunks:
            chunks[-1].extend(chunk)
            break
        chunks.append(chunk)
        i += MAX_BULLETS_PER_SLIDE
    return chunks

# ==============================
# IMAGE HANDLING (UNCHANGED)
# ==============================
def should_use_images(section: str) -> bool:
    return section in {"Method", "Results", "Experiments"}

def page_number_from_path(path: str) -> int:
    m = re.search(r'page_(\d+)', os.path.basename(path))
    return int(m.group(1)) if m else 0

def select_best_images(images, section, used_images, max_images=1):
    candidates = [p for p in images if p not in used_images]
    if not candidates:
        return []
    candidates.sort(key=page_number_from_path)
    return candidates[:max_images] if section == "Method" else candidates[-max_images:]

def generate_image_caption(section: str) -> str:
    if section == "Method":
        return "Model architecture overview."
    if section in {"Results", "Experiments"}:
        return "Experimental results and performance comparison."
    return ""

# ==============================
# STEP-4: DEDUPLICATION & DENSITY CONTROL
# ==============================

# SECTION_TARGET_BULLETS used from top of file

def deduplicate_bullets(bullets, overlap_threshold=0.7):
    """
    Removes near-duplicate bullets using token overlap.
    SAFE: extractive only.
    """
    def normalize(b):
        b = b.lower()
        b = re.sub(r'[^a-z\s]', '', b)
        return set(t for t in b.split() if len(t) > 3)

    kept = []
    seen = []

    for b in bullets:
        tb = normalize(b)
        duplicate = False
        for prev in seen:
            if not tb or not prev:
                continue
            overlap = len(tb & prev) / max(len(tb), len(prev))
            if overlap >= overlap_threshold:
                duplicate = True
                break
        if not duplicate:
            kept.append(b)
            seen.append(tb)

    return kept


_LOW_SIGNAL_PATTERNS = [
    r"hyperparameters",
    r"training steps",
    r"batch size",
    r"learning rate",
    r"byte[- ]pair encoding",
    r"wordpiece",
    r"unlisted values",
]

def remove_low_signal_bullets(bullets):
    clean = []
    for b in bullets:
        low = b.lower()
        if any(re.search(p, low) for p in _LOW_SIGNAL_PATTERNS):
            continue
        clean.append(b)
    return clean


# ==============================
# QWEN SUMMARIZATION LOGIC
# ==============================

def extract_title_with_qwen(first_page_text: str) -> str:
    if not HAS_LLM:
        return ""
    prompt = f"""
    Identify the title of the research paper from the following text (first page content).
    Return ONLY the title, nothing else. Do not add quotes.
    
    TEXT:
    {first_page_text[:2000]}
    
    TITLE:
    """
    try:
        title = qwen_generate(prompt, max_tokens=64).strip()
        # Clean up if it's too long or has newlines
        title = title.replace("\n", " ")
        # Remove common prefixes if Qwen generated them
        if title.lower().startswith("title:"):
            title = title[6:].strip()
        if len(title) > 300: # Sanity check
             return "" 
        return title
    except Exception as e:
        print(f"[WARN] Title extraction failed: {e}")
        return ""

def summarize_section_with_qwen(section_title: str, section_text: str, target_bullets: int) -> List[str]:
    """
    Uses Qwen to generate high-quality bullet points for the slide.
    """
    if not HAS_LLM:
        return []

    print(f"[INFO] Using Qwen LLM for summarization of section: {section_title}")

    prompt = f"""
You are an expert research scientist assisting in creating a high-quality presentation.
Your goal is to EXTRACT key technical details from the provided text into clear, standalone bullet points.
Avoid generic summaries. Use specific details, numbers, and terminology from the text.

SECTION: {section_title}
TEXT:
{section_text[:4000]}

INSTRUCTIONS:
1. Extract exactly {target_bullets + 2} distinct key points.
2. Each bullet must be a COMPLETE sentence ending with a period.
3. BE SPECIFIC: Include metrics, method names, and architectural details if present.
4. MATH & TABLES: 
   - If meaningful tables exist, extract their key insights as bullets.
   - If important mathematical formulas exist, include them as readable text or simple LaTeX (e.g., "Loss L = ...") within a sentence.
   - Ensure these are integrated naturally as bullet points.
5. NO FILLER: Do not use "The paper discusses...", "This section shows...", etc. Start directly with the fact.
6. NO HALLUCINATIONS: Only usage information present in the text.
7. IGNORE citations (e.g. [12]), figures (e.g. Fig 1), and acknowledgments.
8. FORMAT: Return a simple list where each line starts with "- ".

OUTPUT:
"""
    try:
        response = qwen_generate(prompt, max_tokens=1024)
        # Parse bullets
        bullets = []
        for line in response.split('\n'):
            line = line.strip()
            # Strict parsing for bullets
            if line.startswith("- ") or line.startswith("* "):
                clean_line = line[2:].strip()
                if len(clean_line) > 15: # slightly stricter length check
                     # Ensure it looks like a sentence
                    if not clean_line.endswith('.'):
                        clean_line += '.'
                    bullets.append(f"* {clean_line}")
            
        return bullets
    except Exception as e:
        print(f"[WARNING] Qwen summarization failed for {section_title}: {e}")
        return []


def limit_section_bullets(section, bullets):
    target = SECTION_TARGET_BULLETS.get(section)
    if not target:
        return bullets
    return bullets[:target]


# ==============================
# CORE PIPELINE (STABLE)
# ==============================
def generate_slides(input_pdf: str, output_ppt: str, max_bullets=4):
    pages_text, pages_images = load_input_paper(input_pdf)
    sections = split_into_sections(pages_text)

    slides_plan = []
    used_images = set()

    for sec in sections:
        section = normalize_section(sec.get("raw_title") or sec.get("title") or "")
        if section.lower() in SKIP_SECTIONS:
            continue

        text = normalize_pdf_text(sec.get("text", ""))[:MAX_MODEL_CHARS]
        sentences = extract_sentences(text)

        target = SECTION_TARGET_BULLETS.get(section, 15)
        
        # --- NEW: TRY QWEN FIRST ---
        bullets = summarize_section_with_qwen(section, text, target)
        
        # --- FALLBACK: USE RULE-BASED IF QWEN FAILS OR RETURNS NOTHING ---
        if not bullets:
            print(f"[INFO] Qwen yielded no bullets for {section}, using regex fallback.")
            rewritten = [rewrite_bullet(s) for s in sentences]
            bullets = final_bullets(rewritten)
            
            rewritten_all = [rewrite_bullet(s) for s in sentences]
            bullets = add_one_more_safe_bullet(bullets, sentences, target)
            bullets = [light_polish_bullet(b) for b in bullets]
            bullets = enhance_bullet_flow(bullets)
            bullets = polish_bullets(bullets)
            bullets = deduplicate_bullets(bullets)
            bullets = remove_low_signal_bullets(bullets)
            bullets = limit_section_bullets(section, bullets)
        
        # If Qwen worked, we still run a light cleanup pass
        else:
             bullets = [light_polish_bullet(b) for b in bullets]
             # Ensure we don't have too many if Qwen hallucinated extra
             if len(bullets) > target + 5:
                 bullets = bullets[:target+5]

        bullets = [light_polish_bullet(b) for b in bullets]


        if not bullets:
            continue

        bullets = enhance_bullet_flow(bullets)
        bullets = polish_bullets(bullets)

        # ===== DEDUPLICATION & FILTERING =====
        bullets = deduplicate_bullets(bullets)
        bullets = remove_low_signal_bullets(bullets)
        bullets = limit_section_bullets(section, bullets)
        # =====================================

        bullet_chunks = chunk_bullets(bullets)



        images = []
        if should_use_images(section):
            all_paths = [
                img["path"]
                for imgs in pages_images.values()
                for img in imgs
                if "path" in img
            ]
            selected = select_best_images(
                all_paths, section, used_images, MAX_FIGURES_PER_SLIDE
            )
            images = [{"path": p, "caption": generate_image_caption(section)} for p in selected]
            used_images.update(selected)

        for i, chunk in enumerate(bullet_chunks):
            slides_plan.append({
                "title": section if i == 0 else f"{section} (continued)",
                "bullets": chunk,
                "images": images if i == 0 else []
            })




    # --- INTELLIGENT TITLE EXTRACTION ---
    print("[paper2ppt] Determining presentation title...")
    doc_title = ""
    if pages_text:
        # 1. Try Qwen Extraction
        doc_title = extract_title_with_qwen(pages_text[0])
        if doc_title:
             print(f"[paper2ppt] Title extracted via Qwen: {doc_title}")
        
        # 2. Fallback to first line
        if not doc_title or len(doc_title) < 5:
            doc_title = pages_text[0].split("\n")[0]
            print(f"[paper2ppt] Title extracted via First Line: {doc_title}")

    # 3. Last resort fallback
    if not doc_title or len(doc_title) < 5:
        doc_title = Path(input_pdf).stem
        print(f"[paper2ppt] Title fallback to filename: {doc_title}")
    
    # -------------------------------------

    ppt = build_presentation(slides_plan, output_ppt, doc_title, sections)
    return ppt, slides_plan


def attach_narration(presentation, slides_plan):
    for slide, plan in zip(presentation.slides, slides_plan):
        narration = generate_explanatory_narration(
            slide_title=plan["title"],
            bullets=plan["bullets"],
            section=plan["title"],   # section == normalized title
            has_image=bool(plan["images"]),
        )
        slide.notes_slide.notes_text_frame.text = narration


from typing import List

def generate_explanatory_narration(
    slide_title: str,
    bullets: List[str],
    section: str,
    has_image: bool,
) -> str:

    """
    Generate explanatory narration that clarifies the slide
    without adding new information.
    """

    narration = []

    # 1. Context setting
    narration.append(
        f"This slide is part of the {section.lower()} section and focuses on {slide_title.lower()}."
    )

    # Optional image mention (descriptive only)
    if has_image:
        narration.append(
            "The visual on this slide helps illustrate the main idea being described."
        )

    # 2. Bullet-by-bullet explanation
    for b in bullets:
        text = b[2:] if b.startswith("* ") else b
        narration.append(
            f"Here, the slide explains that {text[0].lower() + text[1:]}"
        )

    # 3. Light synthesis
    narration.append(
        "Taken together, these points clarify the key idea presented on this slide."
    )

    # 4. Transition cue
    narration.append(
        "This prepares the ground for the discussion that follows."
        if section.lower() not in {"conclusion"}
        else
        "This concludes the main message of the paper."
    )

    return " ".join(narration)


# ==============================
# CLI
# ==============================
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("-i", "--input", required=True)
    ap.add_argument("-o", "--output", required=True)
    ap.add_argument("--max-bullets", type=int, default=4)
    ap.add_argument(
    "--narration",
    choices=["explanatory"],
    help="Generate hidden narration (speaker notes)"
)

    args = ap.parse_args()

    # First generate the slides (returns the output path)
    out, slides_plan = generate_slides(args.input, args.output, args.max_bullets)

    if args.narration == "explanatory":
        prs = Presentation(args.output)

        for slide, plan in zip(prs.slides, slides_plan):
            slide_text = "\n".join(
                b[2:] if b.startswith("* ") else b
                for b in plan["bullets"]
            )

            narration = generate_narration(
                title=plan["title"],
                slide_text=slide_text,
                summary=plan["title"],
            )

            if narration:
                slide.notes_slide.notes_text_frame.text = narration

        prs.save(args.output)

    print("Saved:", out)




if __name__ == "__main__":
    main()
