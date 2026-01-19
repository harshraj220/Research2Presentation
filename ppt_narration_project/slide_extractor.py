from pptx import Presentation # type: ignore

def extract_slides(ppt_path: str):
    prs = Presentation(ppt_path)
    slides = []

    for slide in prs.slides:
        title = ""
        body_text = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if not text:
                    continue
                if not title:
                    title = text
                else:
                    body_text.append(text)

        slides.append({
            "slide_title": title,
            "original_slide_text": "\n".join(body_text)
        })

    return slides
