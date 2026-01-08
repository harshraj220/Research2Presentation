from pptx import Presentation
from pptx.util import Inches

def embed_audio(input_ppt, audio_dir, output_ppt):
    prs = Presentation(input_ppt)

    for idx, slide in enumerate(prs.slides, start=1):
        audio_path = f"{audio_dir}/slide_{idx}.wav"

        try:
            left = Inches(0)
            top = Inches(0)
            width = Inches(1)
            height = Inches(1)

            slide.shapes.add_movie(
                audio_path,
                left,
                top,
                width,
                height,
                poster_frame_image=None,
                mime_type="audio/wav"
            )
        except FileNotFoundError:
            print(f"[WARN] Audio not found for slide {idx}")

    prs.save(output_ppt)
