from pptx import Presentation

def add_speaker_notes(input_ppt_path, narrations, output_ppt_path):
    prs = Presentation(input_ppt_path)

    for idx, slide in enumerate(prs.slides):
        if idx >= len(narrations):
            break

        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = narrations[idx]


    prs.save(output_ppt_path)
