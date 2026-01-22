# =========================================================
# CLEAN WORKING VERSION — macOS-compatible offline TTS
# =========================================================

import os, re, traceback, sys, subprocess, shutil
from datetime import datetime
from io import BytesIO
from typing import List, Dict, Tuple
from pathlib import Path
import fitz
from PIL import Image, ImageOps
from tqdm import tqdm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Optional imports for summarizer
try:
    from transformers import pipeline
except Exception:
    pipeline = None

# Offline TTS
try:
    import pyttsx3
except Exception:
    pyttsx3 = None

try:
    from gtts import gTTS
except Exception:
    gTTS = None

# ============ CONFIG ============

SUMMARY_MODEL = None         # Disable transformer model
TARGET_BULLETS_PER_SECTION = 20  # Increased to allow more content

BULLETS_PER_SLIDE = 5
MAX_SECTION_CHARS = 4000
AUDIO_OUT_DIR = Path("paper2ppt_audio")
AUDIO_OUT_DIR.mkdir(exist_ok=True)

# ============ CLEANING FUNCTIONS ============

def clean_academic_noise(t: str) -> str:
    if not t:
        return ""
    t = re.sub(r"\[\d+\]", " ", t)
    t = re.sub(r"https?://\S+|www\.\S+", " ", t)
    t = re.sub(r"\s+", " ", t)
    return t.strip()

# Minimal heading detection
def is_heading_line(line: str) -> bool:
    line = line.strip()
    if len(line) > 120:
        return False
    if line.lower() in ("abstract", "introduction"):
        return True
    if re.match(r"^\d+(\.\d+)*\s+[A-Za-z]", line):
        return True
    return False

def normalize_heading(t: str) -> str:
    t = t.lower().strip()
    if "abstract" in t:
        return "abstract"
    if "intro" in t:
        return "introduction"
    return t

# ============ PDF LOADING ============

def read_pdf(path: str):
    doc = fitz.open(path)
    pages = []
    page_images = {}
    for i, p in enumerate(doc):
        pages.append(p.get_text())
        imgs = []
        for img in p.get_images(full=True):
            try:
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                if pix.n >= 5:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                out = f"/tmp/page_{i+1}_img.png"
                pix.save(out)
                imgs.append(out)
            except Exception:
                pass
        page_images[i] = imgs
    return pages, page_images

# ============ SECTION SPLITTING ============

def split_into_sections(pages):
    sections = []
    current = None
    for i, text in enumerate(pages):
        lines = text.splitlines()
        for ln in lines:
            if is_heading_line(ln):
                if current:
                    sections.append(current)
                current = {
                    "title": normalize_heading(ln),
                    "text": "",
                    "pages": {i}
                }
            else:
                if current is None:
                    current = {"title":"title","text":"", "pages":{i}}
                current["text"] += " " + ln
    if current:
        sections.append(current)
    # clean
    for s in sections:
        s["text"] = clean_academic_noise(s["text"])
    return sections

# ============ SUMMARIZATION (HEURISTIC) ============

def summarize_section_to_bullets(text: str, n: int = 5) -> List[str]:
    """
    No transformer — heuristic split.
    Preserves punctuation to ensure sentences look complete.
    """
    # Split keeping the delimiter
    parts = re.split(r"([.!?])", text)
    
    bullets = []
    current_sentence = ""
    
    for p in parts:
        if not p:
            continue
        if p in ".!?":
            current_sentence += p
            if len(current_sentence.strip()) > 10:  # slight increase in min length
                bullets.append(current_sentence.strip())
            current_sentence = ""
        else:
            current_sentence += p
            
    # Catch any remainder
    if len(current_sentence.strip()) > 10:
        bullets.append(current_sentence.strip())

    return bullets[:n] if bullets else ["Summary not available."]

# ============ NARRATION ============

def generate_narration_from_bullets(bullets: List[str], max_words=150):
    t = " ".join(bullets)
    narration = "Here is a brief explanation of this slide: " + t
    # Instead of hard cutting words, we try to keep it reasonable but complete
    words = narration.split()
    if len(words) <= max_words:
        return narration
    
    # If we must cut, try to cut at a punctuation
    truncated = " ".join(words[:max_words])
    # finding last punctuation
    last_dot = truncated.rfind('.')
    last_excl = truncated.rfind('!')
    last_q = truncated.rfind('?')
    cut_idx = max(last_dot, last_excl, last_q)
    
    if cut_idx > 0:
        return truncated[:cut_idx+1]
    
    return truncated + "..."

# ============ TTS ENGINE ============

def synthesize_audio(narration_text: str, slide_idx: int):
    """
    Try pyttsx3 → macOS say → gTTS
    """
    base_mp3 = AUDIO_OUT_DIR / f"slide_{slide_idx}.mp3"

    # ----- 1) pyttsx3 -----
    if pyttsx3 is not None:
        try:
            eng = pyttsx3.init()
            wav_path = AUDIO_OUT_DIR / f"slide_{slide_idx}.wav"
            eng.save_to_file(narration_text, str(wav_path))
            eng.runAndWait()
            # convert to mp3 if ffmpeg exists
            if shutil.which("ffmpeg"):
                subprocess.run(["ffmpeg","-y","-i",str(wav_path),str(base_mp3)],
                               stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                wav_path.unlink(missing_ok=True)
            return base_mp3
        except Exception as e:
            print(f"[Slide {slide_idx}] pyttsx3 failed: {e}")

    # ----- 2) macOS say -----
    if sys.platform == "darwin":
        try:
            aiff = AUDIO_OUT_DIR / f"slide_{slide_idx}.aiff"
            subprocess.run(["say","-o",str(aiff),narration_text], check=True)
            if shutil.which("ffmpeg"):
                subprocess.run(["ffmpeg","-y","-i",str(aiff),str(base_mp3)],
                               stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                aiff.unlink(missing_ok=True)
                return base_mp3
            return aiff
        except Exception as e:
            print(f"[Slide {slide_idx}] say failed: {e}")

    # ----- 3) gTTS fallback -----
    if gTTS is not None:
        try:
            tts = gTTS(narration_text)
            tts.save(str(base_mp3))
            return base_mp3
        except Exception as e:
            print(f"[Slide {slide_idx}] gTTS failed: {e}")

    return None

# ============ PPT CREATION ============

def create_ppt_with_audio(slides_plan, output_name, pages_text, infile):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # Title slide
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = pages_text[0].split("\n")[0][:100]
    try:
        sl.placeholders[1].text = f"Auto-generated • {datetime.now().strftime('%Y-%m-%d')}"
    except:
        pass

    # Content slides
    slide_idx = 0
    created_audio = []

    for s in tqdm(slides_plan, desc="Composing slides"):
        slide_idx += 1
        slide = prs.slides.add_slide(blank)

        # Panel
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0),
                                       prs.slide_width - Inches(1.6),
                                       prs.slide_height - Inches(2))
        tf = box.text_frame

        # Title
        p = tf.add_paragraph()
        p.text = s["title"]
        p.font.size = Pt(28)
        p.font.bold = True

        # Bullets
        for b in s["bullets"]:
            pb = tf.add_paragraph()
            pb.text = "• " + b
            pb.level = 1
            pb.font.size = Pt(18)

        # Narration
        narration = generate_narration_from_bullets(s["bullets"])
        print(f"[Slide {slide_idx}] narration preview:", narration[:80])

        # Notes
        try:
            slide.notes_slide.notes_text_frame.text = narration
        except:
            pass

        # TTS
        audio_path = synthesize_audio(narration, slide_idx)
        if audio_path:
            print(f"[Slide {slide_idx}] saved audio:", audio_path.name)
            created_audio.append(str(audio_path))
        else:
            print(f"[Slide {slide_idx}] WARNING no audio created")

    prs.save(output_name)
    return output_name, created_audio

# ============ MAIN ============

def main(input_file, output_file):
    ext = input_file.split(".")[-1].lower()

    if ext == "pdf":
        pages_text, pages_images = read_pdf(input_file)
    else:
        pages_text = [open(input_file).read()]
        pages_images = {0:[]}

    sections = split_into_sections(pages_text)

    slides_plan = []
    
    for s in sections:
        # 1. Get ALL valid bullets (up to limit)
        section_bullets = summarize_section_to_bullets(s["text"], TARGET_BULLETS_PER_SECTION)
        
        # 2. Chunk them for pagination
        chunk_size = BULLETS_PER_SLIDE
        if not section_bullets:
             # If empty, maybe just skip or add a placeholder?
             # Let's verify if we want empty sections.
             continue
             
        # Create chunks
        # e.g. [0, 5, 10...]
        for i in range(0, len(section_bullets), chunk_size):
            chunk = section_bullets[i : i + chunk_size]
            
            # Determine title
            slide_title = s["title"].title()
            if i > 0:
                slide_title += " (Continued)"
            
            slides_plan.append({
                "title": slide_title,
                "bullets": chunk
            })

    if not slides_plan:
        raise RuntimeError("No slides produced.")

    out, aud = create_ppt_with_audio(slides_plan, output_file, pages_text, input_file)
    print("DONE →", out)
    print("Audio files:", aud)


if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("-i", "--input", required=True)
    ap.add_argument("-o", "--output", required=True)
    args = ap.parse_args()
    main(args.input, args.output)

