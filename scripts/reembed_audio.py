import os
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

PPTX = "test_output.pptx"
AUDIO_DIR = Path("paper2ppt_audio")

prs = Presentation(PPTX)
total = len(prs.slides)
print(f"Opening {PPTX} with {total} slides. Looking for audio files in {AUDIO_DIR}/")

embedded = []
failed = []

for i, slide in enumerate(prs.slides, start=1):
    audio_path = AUDIO_DIR / f"slide_{i}.mp3"
    if not audio_path.exists():
        alt_path = AUDIO_DIR / f"slide_{i-1}.mp3"
        if alt_path.exists():
            audio_path = alt_path
        else:
            print(f"[Slide {i}] No audio file found ({audio_path}) — skipping.")
            continue

    try:
        left = Inches(0.2); top = Inches(0.2); width = Inches(0.8); height = Inches(0.8)
        slide.shapes.add_movie(str(audio_path), left, top, width, height, mime_type="audio/mpeg")
        embedded.append((i, str(audio_path)))
        print(f"[Slide {i}] Embedded {audio_path.name}")
    except Exception as e:
        failed.append((i, str(audio_path), str(e)))
        print(f"[Slide {i}] FAILED to embed {audio_path.name}: {e}")

outname = "test_output_with_audio_embedded.pptx"
prs.save(outname)
print("Done. Embedded on", len(embedded), "slides; failed on", len(failed))
print("Saved new PPTX as", outname)

if failed:
    print("Failures (first 5):")
    for f in failed[:5]:
        print(f)
