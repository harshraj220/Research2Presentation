from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pathlib import Path
from typing import List, Dict
from datetime import datetime
from PIL import Image, ImageOps, ImageChops
from io import BytesIO
import os

# =========================
# THEME & CONSTANTS
# =========================

DEFAULT_THEME = {
    "title_font": "Calibri",
    "body_font": "Calibri",
    "title_size_pt": 32,
    "body_size_pt": 16,
    "accent_color": [18, 90, 173],
    "background_color": [255, 255, 255],
    "panel_bg": [255, 255, 255],
    "text_color": [0, 0, 0],
    "subtext_color": [90, 90, 90],
    "panel_transparency": 0.0,
}

MAX_FIGURES_PER_SLIDE = 2
MAX_VISIBLE_BULLETS = 4
MAX_BULLET_CHARS = 140

# =========================
# UTILS
# =========================

def rgb(c):
    return RGBColor(int(c[0]), int(c[1]), int(c[2]))


def _draw_footer(slide, prs, theme):
    footer = slide.shapes.add_textbox(
        Inches(0.6),
        prs.slide_height - Inches(0.45),
        prs.slide_width - Inches(1.2),
        Inches(0.3),
    )
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"Auto-generated • {datetime.utcnow().strftime('%Y-%m-%d')}"
    p.font.size = Pt(9)
    p.font.color.rgb = rgb(theme["subtext_color"])
    p.alignment = PP_ALIGN.CENTER


# =========================
# IMAGE HELPERS
# =========================

from PIL import Image, ImageChops, ImageFile
from io import BytesIO
import os

ImageFile.LOAD_TRUNCATED_IMAGES = True  # ✅ CRITICAL

def crop_image_whitespace(path: str) -> BytesIO:
    """
    Robust image loader.
    - Handles truncated PNGs
    - Falls back gracefully
    - Never crashes slide generation
    """
    try:
        if not os.path.exists(path) or os.path.getsize(path) < 1024:
            raise OSError("Image missing or too small")

        img = Image.open(path)
        img.load()
        img = img.convert("RGB")

        bg = Image.new(img.mode, img.size, img.getpixel((0, 0)))
        diff = ImageChops.difference(img, bg)
        bbox = diff.getbbox()

        if bbox:
            img = img.crop(bbox)

    except Exception as e:
        print(f"[WARN] Skipping corrupted image: {path} ({e})")
        img = Image.new("RGB", (800, 600), (255, 255, 255))

    bio = BytesIO()
    img.save(bio, format="PNG")
    bio.seek(0)
    return bio



def fit_image(bio, w_px, h_px) -> BytesIO:
    """Resize image to fit within dimensions while maintaining aspect ratio"""
    try:
        bio.seek(0)
        img = Image.open(bio)
        
        # Use PIL's contain to maintain aspect ratio
        img = ImageOps.contain(img, (w_px, h_px))
        
        out = BytesIO()
        img.save(out, format="PNG")
        out.seek(0)
        return out
        
    except Exception as e:
        print(f"[WARN] Failed to resize image: {e}")
        bio.seek(0)
        return bio


# =========================
# SLIDE BUILDERS
# =========================

def add_title(slide, prs, text, theme):
    """Add title to slide"""
    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.4),
        prs.slide_width - Inches(1.2), Inches(0.9)
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text[:120]
    p.font.size = Pt(theme["title_size_pt"])
    p.font.bold = True
    p.font.color.rgb = rgb(theme["text_color"])
    p.alignment = PP_ALIGN.LEFT


def add_bullets(slide, prs, bullets, has_images, theme):
    """Add bullet points to slide with proper spacing"""
    left = Inches(0.6)
    top = Inches(1.5)
    
    # Adjust width based on whether we have images
    if has_images:
        width = prs.slide_width - Inches(4.6)  # Leave space for images
    else:
        width = prs.slide_width - Inches(1.2)
    
    height = prs.slide_height - Inches(2.0)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for i, b in enumerate(bullets[:MAX_VISIBLE_BULLETS]):
        if not b:
            continue

        # Safety truncation
        if len(b) > MAX_BULLET_CHARS:
            b = b[:MAX_BULLET_CHARS - 1] + "…"

        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        p.font.size = Pt(theme["body_size_pt"])
        p.font.color.rgb = rgb(theme["text_color"])
        p.space_after = Pt(12)


def add_images(slide, prs, images, theme):
    """Add images to the right side of slide"""
    if not images:
        return
    
    # Image column dimensions
    col_w = Inches(3.6)
    col_h = Inches(2.6)
    gap = Inches(0.4)
    left = prs.slide_width - col_w - Inches(0.5)
    top = Inches(1.4)

    # Convert to pixels for PIL
    px_w = int(col_w.inches * 96)
    px_h = int(col_h.inches * 96)

    added_count = 0
    for i, img in enumerate(images[:MAX_FIGURES_PER_SLIDE]):
        try:
            img_path = img.get("path", "")
            
            # Verify file exists
            if not os.path.exists(img_path):
                print(f"[ERROR] Image not found: {img_path}")
                continue
            
            print(f"[IMAGE] Adding to slide: {os.path.basename(img_path)}")
            
            # Process image
            bio = crop_image_whitespace(img_path)
            bio = fit_image(bio, px_w, px_h)
            
            # Calculate position for multiple images
            img_top = top + i * (col_h + gap)
            
            # Add image to slide
            pic = slide.shapes.add_picture(
                bio, left, img_top, col_w, col_h
            )
            
            added_count += 1
            
            # Add caption if available
            caption_text = img.get("caption", "").strip()
            if caption_text:
                cap = slide.shapes.add_textbox(
                    left,
                    img_top + col_h + Inches(0.05),
                    col_w,
                    Inches(0.35),
                )
                tf = cap.text_frame
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = caption_text[:180]
                p.font.size = Pt(10)
                p.font.italic = True
                p.font.color.rgb = rgb(theme["subtext_color"])
                p.alignment = PP_ALIGN.CENTER
                
        except Exception as e:
            print(f"[ERROR] Failed to add image {img.get('path', 'unknown')}: {e}")
            import traceback
            traceback.print_exc()
            continue
    
    if added_count > 0:
        print(f"[SUCCESS] Added {added_count} images to slide")
    else:
        print(f"[WARN] No images were successfully added to slide")


# =========================
# MAIN BUILDER
# =========================

def build_presentation(slides_plan, output_path, doc_title, sections, theme_name=None):
    """
    Build PowerPoint presentation from slides plan.
    
    Args:
        slides_plan: List of slide dicts with keys: title, bullets, images
        output_path: Path to save PPTX
        doc_title: Document title for first slide
        sections: Original sections (for reference)
        theme_name: Optional theme name
    """
    theme = DEFAULT_THEME
    prs = Presentation()
    blank = prs.slide_layouts[6]  # Blank layout

    print(f"\n[PPTX] Building presentation with {len(slides_plan)} slides...")

    for idx, s in enumerate(slides_plan):
        slide = prs.slides.add_slide(blank)
        
        # Force solid white background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Get slide content
        title = s.get("title", "")
        bullets = s.get("bullets", [])
        images = s.get("images", [])
        
        print(f"\n[SLIDE {idx+1}] Title: '{title}'")
        print(f"  Bullets: {len(bullets)}")
        print(f"  Images: {len(images)}")
        
        # Add content
        add_title(slide, prs, title, theme)
        add_bullets(slide, prs, bullets, bool(images), theme)
        
        if images:
            print(f"  Processing {len(images)} images...")
            add_images(slide, prs, images, theme)
        
        _draw_footer(slide, prs, theme)

    # Save presentation
    prs.save(output_path)
    print(f"\n[SUCCESS] Presentation saved: {output_path}")
    
    return output_path