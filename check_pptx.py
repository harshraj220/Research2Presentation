from pptx import Presentation
import sys

def inspect_pptx(filename):
    try:
        prs = Presentation(filename)
        print(f"--- CONTENT OF {filename} ---")
        for i, slide in enumerate(prs.slides):
            title = slide.shapes.title.text if slide.shapes.title else "NO TITLE"
            print(f"\n[Slide {i+1}] {title}")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape == slide.shapes.title:
                    continue
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        print(f"  - {p.text.strip()}")
    except Exception as e:
        print(f"Error reading pptx: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        inspect_pptx(sys.argv[1])
    else:
        inspect_pptx("output.pptx")
