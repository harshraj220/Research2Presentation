
import fitz
from pptx import Presentation
import os

pdf_path = "doc_to_ppt.pdf"
ppt_path = "doc_to_ppt_summary_with_narration.pptx"

print(f"--- CHECKING INPUT: {pdf_path} ---")
if not os.path.exists(pdf_path):
    print("PDF NOT FOUND")
else:
    try:
        doc = fitz.open(pdf_path)
        pdf_text = ""
        for page in doc:
            pdf_text += page.get_text()
        print(f"PDF extracted {len(pdf_text)} chars.")
        print("First 500 chars of PDF:")
        print(pdf_text[:500])
    except Exception as e:
        print(f"Error reading PDF: {e}")

print(f"\n--- CHECKING OUTPUT: {ppt_path} ---")
if not os.path.exists(ppt_path):
    # Try alternate name if script hasn't finished renaming or failed
    ppt_path = "output.pptx"
    print(f"Target PPTX not found, checking intermediate: {ppt_path}")

if not os.path.exists(ppt_path):
    print("PPTX NOT FOUND")
else:
    try:
        prs = Presentation(ppt_path)
        print(f"PPTX has {len(prs.slides)} slides.")
        
        for i, slide in enumerate(prs.slides):
            print(f"\nSlide {i+1}:")
            # Title
            if slide.shapes.title:
                print(f"Title: {slide.shapes.title.text}")
            
            # Text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text:
                             print(f"Bullet: {paragraph.text}")
                elif hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    print(f"Body: {shape.text}")
                
            # Check notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                print(f"Notes (Narration): {notes}")
    except Exception as e:
        print(f"Error reading PPTX: {e}")
